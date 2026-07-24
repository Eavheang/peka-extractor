"""Extract local source files into JSON records."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

import fitz
import pdfplumber
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation


SECTION_HEADING_RE = re.compile(r"^(?P<number>\d+)[.)]\s+(?P<title>[A-Z][^\n]*)$")

SUPPORTED = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}


def clean_text(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        for chunk in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def source_id_for(path: Path) -> str:
    identity = path.resolve().as_posix().casefold()
    return hashlib.sha256(identity.encode("utf-8")).hexdigest()[:16]


def extract_pdf(path: Path) -> list[dict[str, object]]:
    with fitz.open(path) as document:
        pymupdf_text = [clean_text(page.get_text()) for page in document]

    try:
        with pdfplumber.open(path) as document:
            pdfplumber_text = [
                clean_text(page.extract_text() or "") for page in document.pages
            ]
    except Exception:
        pdfplumber_text = []

    sections = []
    for number in range(max(len(pymupdf_text), len(pdfplumber_text))):
        text = pymupdf_text[number] if number < len(pymupdf_text) else ""
        if not text and number < len(pdfplumber_text):
            text = pdfplumber_text[number]
        if text:
            sections.append({"page": number + 1, "text": text})
    return sections


def extract_docx(path: Path) -> list[dict[str, object]]:
    document = Document(path)
    section_blocks: dict[int, list[str]] = {0: []}
    current_section = 0

    for child in document.element.body.iterchildren():
        if child.tag.endswith("}p"):
            paragraph = Paragraph(child, document)
            raw_text = paragraph.text
            text = clean_text(raw_text)
            style_name = paragraph.style.name if paragraph.style else ""
            match = SECTION_HEADING_RE.fullmatch(text)
            is_heading = match and (
                style_name.lower().startswith("heading")
                or int(match.group("number")) == current_section + 1
            )
            if is_heading:
                current_section = int(match.group("number"))
                section_blocks.setdefault(current_section, [])
        elif child.tag.endswith("}tbl"):
            table = Table(child, document)
            text = clean_text(
                "\n".join(
                    "\t".join(clean_text(cell.text) for cell in row.cells)
                    for row in table.rows
                )
            )
        else:
            continue

        if text:
            section_blocks.setdefault(current_section, []).append(text)

    return [
        {"section": number, "text": clean_text("\n".join(blocks))}
        for number, blocks in section_blocks.items()
        if clean_text("\n".join(blocks))
    ]


def extract_pptx(path: Path) -> list[dict[str, object]]:
    presentation = Presentation(path)
    sections = []
    for number, slide in enumerate(presentation.slides, start=1):
        text = clean_text(
            "\n".join(
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            )
        )
        if text:
            sections.append({"slide": number, "text": text})
    return sections


def extract_xlsx(path: Path) -> list[dict[str, object]]:
    try:
        import pandas
    except ImportError:
        pandas = None

    workbook = load_workbook(path, read_only=True, data_only=True)
    sheet_names = workbook.sheetnames
    sections = []
    for number, sheet_name in enumerate(sheet_names, start=1):
        if pandas is not None:
            frame = pandas.read_excel(
                path,
                sheet_name=sheet_name,
                header=None,
                dtype=object,
                engine="openpyxl",
            )
            rows = [
                "\t".join("" if pandas.isna(value) else str(value) for value in row)
                for row in frame.itertuples(index=False, name=None)
            ]
        else:
            worksheet = workbook[sheet_name]
            rows = [
                "\t".join("" if value is None else str(value) for value in row)
                for row in worksheet.iter_rows(values_only=True)
            ]
        text = clean_text("\n".join(rows))
        if text:
            sections.append({"sheet": number, "name": sheet_name, "text": text})
    workbook.close()
    return sections


def extract_file(
    path: Path,
    *,
    source_id: str | None = None,
    content_hash: str | None = None,
) -> dict[str, object]:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED:
        raise ValueError(f"Unsupported file type: {path.suffix or '(none)'}")

    resolved_path = path.resolve()
    source_sha256 = content_hash or file_sha256(resolved_path)
    if suffix == ".pdf":
        sections = extract_pdf(resolved_path)
    elif suffix == ".docx":
        sections = extract_docx(resolved_path)
    elif suffix == ".pptx":
        sections = extract_pptx(resolved_path)
    elif suffix == ".xlsx":
        sections = extract_xlsx(resolved_path)
    else:
        text = clean_text(resolved_path.read_text(encoding="utf-8", errors="replace"))
        sections = [{"text": text}] if text else []

    return {
        "source_id": source_id or source_id_for(resolved_path),
        "source_sha256": source_sha256,
        "source_filename": resolved_path.name,
        "source_path": str(resolved_path),
        "source_extension": suffix,
        "extracted_at": datetime.now(UTC).isoformat(),
        "sections": sections,
    }


def load_registry(path: Path) -> dict[str, dict[str, Any]]:
    if not path.exists():
        return {}
    raw = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(raw, dict):
        raise ValueError(f"Registry must contain a JSON object: {path}")

    registry = {}
    for key, value in raw.items():
        if not isinstance(value, dict):
            continue
        source_path = value.get("source_path", key)
        if not isinstance(source_path, str):
            continue
        source_id = source_id_for(Path(source_path))
        registry[source_id] = {
            "source_id": source_id,
            "source_sha256": value.get("source_sha256"),
            "source_path": source_path,
            "output": value.get("output"),
        }
    return registry


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract files for LLM processing")
    parser.add_argument("files", type=Path, nargs="+", help="Source files")
    parser.add_argument("--output", type=Path, default=Path("extracted"))
    args = parser.parse_args()
    args.output.mkdir(parents=True, exist_ok=True)
    registry_path = args.output / "registry.json"
    registry = load_registry(registry_path)

    for path in args.files:
        if not path.is_file():
            raise SystemExit(f"File not found: {path}")
        resolved_path = path.resolve()
        source_id = source_id_for(resolved_path)
        source_sha256 = file_sha256(resolved_path)
        previous = registry.get(source_id)
        if previous and previous.get("source_sha256") == source_sha256:
            print(f"Skipped unchanged: {path}")
            continue

        record = extract_file(
            resolved_path,
            source_id=source_id,
            content_hash=source_sha256,
        )
        output_file = args.output / f"{source_id}.json"
        output_file.write_text(json.dumps(record, indent=2) + "\n", encoding="utf-8")
        registry[source_id] = {
            "source_id": source_id,
            "source_sha256": source_sha256,
            "source_path": str(resolved_path),
            "output": str(output_file),
        }
        action = "Updated" if previous else "Extracted"
        print(f"{action}: {path} -> {output_file}")

    registry_path.write_text(json.dumps(registry, indent=2) + "\n", encoding="utf-8")


if __name__ == "__main__":
    main()
